# -*- coding: utf-8 -*-
"""
구글 노트북LM 이미지 슬라이드를 편집 가능한 PPTX로 변환하는 앱
PyQt5 + EasyOCR + python-pptx
"""

import sys
import os
import io
import fitz  # PyMuPDF
import numpy as np
from PIL import Image
import cv2
from pathlib import Path

from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QLabel, QPushButton, QFileDialog, QLineEdit, QComboBox, QProgressBar,
    QMessageBox, QGroupBox, QFormLayout
)
from PyQt5.QtCore import Qt, QThread, pyqtSignal
from PyQt5.QtGui import QFont

# PPT 16:9 크기 (inches)
PPT_WIDTH = 10.0
PPT_HEIGHT = 5.625


def get_background_color(img: np.ndarray, bbox: np.ndarray, margin: int = 8) -> tuple:
    """텍스트 영역 바깥에서 배경색 추정 (bbox 확장된 모서리에서 샘플링)"""
    h, w = img.shape[:2]
    x_coords = bbox[:, 0]
    y_coords = bbox[:, 1]
    x_min = int(np.min(x_coords)) - margin
    x_max = int(np.max(x_coords)) + margin
    y_min = int(np.min(y_coords)) - margin
    y_max = int(np.max(y_coords)) + margin

    # bbox 바깥 4모서리에서 샘플링 (텍스트가 없는 영역)
    samples = []
    for (sx, sy) in [
        (x_min, y_min), (x_max, y_min), (x_min, y_max), (x_max, y_max),
        (x_min, (y_min + y_max) // 2), (x_max, (y_min + y_max) // 2),
        ((x_min + x_max) // 2, y_min), ((x_min + x_max) // 2, y_max),
    ]:
        sx = np.clip(sx, 0, w - 1)
        sy = np.clip(sy, 0, h - 1)
        samples.append(img[int(sy), int(sx)])

    if not samples:
        for (px, py) in [(5, 5), (w - 6, 5), (5, h - 6), (w - 6, h - 6)]:
            if 0 <= px < w and 0 <= py < h:
                samples.append(img[py, px])

    if not samples:
        return (255, 255, 255)

    arr = np.array(samples)
    return tuple(int(np.median(arr[:, i])) for i in range(min(3, arr.shape[1])))


def mask_text_regions(img: np.ndarray, ocr_results: list) -> np.ndarray:
    """OCR로 감지된 텍스트 영역을 배경색으로 덮어 텍스트 제거"""
    result = img.copy()
    if len(result.shape) == 2:
        result = cv2.cvtColor(result, cv2.COLOR_GRAY2BGR)

    for (bbox, text, conf) in ocr_results:
        if not text.strip():
            continue
        pts = np.array(bbox, dtype=np.int32)
        color = get_background_color(img, np.array(bbox))
        if len(color) == 2:
            color = (color[0], color[0], color[0])
        cv2.fillPoly(result, [pts], color)
    return result


def px_to_inches(px_val: float, img_dim: float, ppt_dim: float) -> float:
    """픽셀 좌표를 PPT 인치로 변환 (비율 기반)"""
    return (px_val / img_dim) * ppt_dim


class ConvertWorker(QThread):
    """변환 작업 스레드"""
    progress = pyqtSignal(int, int, str)
    finished_signal = pyqtSignal(bool, str)

    def __init__(self, pdf_path: str, save_path: str, font_name: str):
        super().__init__()
        self.pdf_path = pdf_path
        self.save_path = save_path
        self.font_name = font_name

    def run(self):
        try:
            import easyocr
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
        except ImportError as e:
            self.finished_signal.emit(False, f"필수 패키지 오류: {e}")
            return

        try:
            self.progress.emit(0, 0, "EasyOCR 초기화 중...")
            reader = easyocr.Reader(["ko", "en"], gpu=False)

            self.progress.emit(0, 0, "PDF 열기...")
            doc = fitz.open(self.pdf_path)
            total = len(doc)
            prs = Presentation()
            prs.slide_width = Inches(PPT_WIDTH)
            prs.slide_height = Inches(PPT_HEIGHT)

            blank_layout = prs.slide_layouts.get_by_name("Blank", prs.slide_layouts[6])

            for page_num in range(total):
                self.progress.emit(page_num + 1, total, f"슬라이드 {page_num + 1}/{total} 처리 중...")

                page = doc[page_num]
                mat = fitz.Matrix(2, 2)  # 2x 해상도로 렌더링
                pix = page.get_pixmap(matrix=mat, alpha=False)
                img_data = pix.tobytes("png")
                img = np.frombuffer(img_data, dtype=np.uint8)
                img = cv2.imdecode(img, cv2.IMREAD_COLOR)
                if img is None:
                    img = np.array(Image.open(io.BytesIO(img_data)))
                    if len(img.shape) == 2:
                        img = cv2.cvtColor(img, cv2.COLOR_GRAY2BGR)

                img_h, img_w = img.shape[:2]

                # OCR
                ocr_results = reader.readtext(img)

                # 텍스트 영역 마스킹 (배경색으로 덮기)
                masked_img = mask_text_regions(img, ocr_results)

                # 메모리 스트림으로 이미지 전달 (임시 파일 사용 시 파일 잠금 오류 방지)
                _, img_buffer = cv2.imencode(".png", masked_img)
                img_stream = io.BytesIO(img_buffer.tobytes())
                img_stream.seek(0)

                # 슬라이드 추가
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(img_stream, Inches(0), Inches(0), Inches(PPT_WIDTH), Inches(PPT_HEIGHT))

                # 추출된 텍스트를 원래 위치에 삽입
                for (bbox, text, conf) in ocr_results:
                    if not text.strip():
                        continue
                    pts = np.array(bbox, dtype=np.float64)
                    x_min, x_max = np.min(pts[:, 0]), np.max(pts[:, 0])
                    y_min, y_max = np.min(pts[:, 1]), np.max(pts[:, 1])
                    w_px = x_max - x_min
                    h_px = y_max - y_min

                    left_in = px_to_inches(x_min, img_w, PPT_WIDTH)
                    top_in = px_to_inches(y_min, img_h, PPT_HEIGHT)
                    width_in = px_to_inches(w_px, img_w, PPT_WIDTH)
                    height_in = px_to_inches(h_px, img_h, PPT_HEIGHT)

                    # 최소 크기
                    width_in = max(width_in, 0.1)
                    height_in = max(height_in, 0.05)

                    txbox = slide.shapes.add_textbox(
                        Inches(left_in), Inches(top_in),
                        Inches(width_in), Inches(height_in)
                    )
                    txbox.fill.background()
                    txbox.line.fill.background()
                    tf = txbox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = text
                    p.font.name = self.font_name
                    p.font.size = Pt(max(6, height_in * 72 / PPT_HEIGHT * 0.8))
                    p.font.color.rgb = RGBColor(0, 0, 0)

            doc.close()
            prs.save(self.save_path)
            self.finished_signal.emit(True, "변환 완료!")

        except Exception as e:
            import traceback
            self.finished_signal.emit(False, str(e) + "\n" + traceback.format_exc())


class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("노트북LM → PPTX 변환기")
        self.setMinimumSize(500, 350)
        self.worker = None
        self.setup_ui()

    def setup_ui(self):
        central = QWidget()
        self.setCentralWidget(central)
        layout = QVBoxLayout(central)

        # PDF 파일
        pdf_group = QGroupBox("PDF 파일")
        pdf_layout = QHBoxLayout()
        self.pdf_edit = QLineEdit()
        self.pdf_edit.setPlaceholderText("PDF 파일 경로를 선택하세요")
        pdf_btn = QPushButton("PDF 불러오기")
        pdf_btn.clicked.connect(self.load_pdf)
        pdf_layout.addWidget(self.pdf_edit)
        pdf_layout.addWidget(pdf_btn)
        pdf_group.setLayout(pdf_layout)
        layout.addWidget(pdf_group)

        # 저장 경로
        save_group = QGroupBox("저장 경로")
        save_layout = QHBoxLayout()
        self.save_edit = QLineEdit()
        self.save_edit.setPlaceholderText("저장할 PPTX 파일 경로를 지정하세요")
        save_btn = QPushButton("경로 지정")
        save_btn.clicked.connect(self.choose_save_path)
        save_layout.addWidget(self.save_edit)
        save_layout.addWidget(save_btn)
        save_group.setLayout(save_layout)
        layout.addWidget(save_group)

        # 폰트 선택
        font_group = QGroupBox("폰트 설정")
        font_layout = QFormLayout()
        self.font_combo = QComboBox()
        self.font_combo.addItems(["나눔고딕", "맑은 고딕"])
        font_layout.addRow("폰트:", self.font_combo)
        font_group.setLayout(font_layout)
        layout.addWidget(font_group)

        # 변환 버튼
        self.convert_btn = QPushButton("변환")
        self.convert_btn.setMinimumHeight(40)
        self.convert_btn.clicked.connect(self.start_convert)
        layout.addWidget(self.convert_btn)

        # 진행률
        self.progress_bar = QProgressBar()
        self.progress_bar.setVisible(False)
        layout.addWidget(self.progress_bar)

        self.status_label = QLabel("")
        layout.addWidget(self.status_label)

    def load_pdf(self):
        path, _ = QFileDialog.getOpenFileName(
            self, "PDF 파일 선택", "", "PDF Files (*.pdf)"
        )
        if path:
            self.pdf_edit.setText(path)
            if not self.save_edit.text():
                base = Path(path).stem + "_editable.pptx"
                self.save_edit.setText(str(Path(path).parent / base))

    def choose_save_path(self):
        path, _ = QFileDialog.getSaveFileName(
            self, "저장 경로 지정", "", "PowerPoint (*.pptx)"
        )
        if path:
            if not path.endswith(".pptx"):
                path += ".pptx"
            self.save_edit.setText(path)

    def start_convert(self):
        pdf_path = self.pdf_edit.text().strip()
        save_path = self.save_edit.text().strip()

        if not pdf_path:
            QMessageBox.warning(self, "오류", "PDF 파일을 선택하세요.")
            return
        if not os.path.isfile(pdf_path):
            QMessageBox.warning(self, "오류", "PDF 파일이 존재하지 않습니다.")
            return
        if not save_path:
            QMessageBox.warning(self, "오류", "저장 경로를 지정하세요.")
            return

        save_dir = os.path.dirname(save_path)
        if save_dir and not os.path.isdir(save_dir):
            QMessageBox.warning(self, "오류", "저장 경로의 폴더가 존재하지 않습니다.")
            return

        font_map = {"나눔고딕": "Nanum Gothic", "맑은 고딕": "Malgun Gothic"}
        font_name = font_map.get(self.font_combo.currentText(), "Malgun Gothic")

        self.convert_btn.setEnabled(False)
        self.progress_bar.setVisible(True)
        self.progress_bar.setValue(0)
        self.status_label.setText("변환 중...")

        self.worker = ConvertWorker(pdf_path, save_path, font_name)
        self.worker.progress.connect(self.on_progress)
        self.worker.finished_signal.connect(self.on_finished)
        self.worker.start()

    def on_progress(self, current: int, total: int, msg: str):
        if total > 0:
            self.progress_bar.setMaximum(total)
            self.progress_bar.setValue(current)
        self.status_label.setText(msg)

    def on_finished(self, success: bool, message: str):
        self.convert_btn.setEnabled(True)
        self.progress_bar.setVisible(False)
        self.status_label.setText("")
        if success:
            QMessageBox.information(self, "완료", message)
        else:
            QMessageBox.critical(self, "오류", message)


def main():
    app = QApplication(sys.argv)
    app.setFont(QFont("Malgun Gothic", 9))
    win = MainWindow()
    win.show()
    sys.exit(app.exec_())


if __name__ == "__main__":
    main()
