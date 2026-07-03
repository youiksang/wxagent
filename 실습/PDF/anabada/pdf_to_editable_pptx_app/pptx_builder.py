"""
python-pptx를 사용해 16:9 편집 가능 PPTX를 생성합니다.
"""

from __future__ import annotations

import os
from typing import Callable, List, Optional, Tuple

import cv2
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from font_utils import ResolvedFont, apply_font_to_run
from ocr_utils import OcrBox

PPT_WIDTH_INCH = 10.0
PPT_HEIGHT_INCH = 5.625

MIN_FONT_SIZE_PT = 8
MAX_FONT_SIZE_PT = 44
DEFAULT_FONT_SIZE_SCALE = 0.8


def px_to_inches(
    x: float,
    y: float,
    w: float,
    h: float,
    image_width: int,
    image_height: int,
    ppt_width: float = PPT_WIDTH_INCH,
    ppt_height: float = PPT_HEIGHT_INCH,
) -> Tuple[float, float, float, float]:
    """
    이미지 픽셀 좌표를 PPT inch 좌표로 변환합니다.

    ppt_x = image_x / image_width_px * ppt_width_inch
    ppt_y = image_y / image_height_px * ppt_height_inch
    ppt_w = image_w / image_width_px * ppt_width_inch
    ppt_h = image_h / image_height_px * ppt_height_inch
    """
    if image_width <= 0 or image_height <= 0:
        raise ValueError("이미지 크기가 유효하지 않습니다.")

    ppt_x = x / image_width * ppt_width
    ppt_y = y / image_height * ppt_height
    ppt_w = w / image_width * ppt_width
    ppt_h = h / image_height * ppt_height

    return ppt_x, ppt_y, ppt_w, ppt_h


def estimate_font_size_pt(
    ppt_h: float,
    scale: float = DEFAULT_FONT_SIZE_SCALE,
    min_pt: int = MIN_FONT_SIZE_PT,
    max_pt: int = MAX_FONT_SIZE_PT,
) -> float:
    """OCR bounding box 높이(inch)를 기준으로 폰트 크기(pt)를 추정합니다."""
    raw = ppt_h * 72.0 * scale
    return float(max(min_pt, min(max_pt, raw)))


def estimate_text_color(
    image_path: str,
    box: OcrBox,
    default_color: Tuple[int, int, int] = (0, 0, 0),
) -> RGBColor:
    """
    OCR 영역 픽셀을 분석해 텍스트 색상을 추정합니다.

    배경과 대비가 큰 색상을 후보로 선택하고, 추정이 어려우면 검정색을 사용합니다.
    """
    image = cv2.imread(image_path)
    if image is None:
        return RGBColor(*default_color)

    height, width = image.shape[:2]
    x1 = max(0, box.x)
    y1 = max(0, box.y)
    x2 = min(width, box.x2)
    y2 = min(height, box.y2)

    if x2 <= x1 or y2 <= y1:
        return RGBColor(*default_color)

    roi = image[y1:y2, x1:x2]
    if roi.size == 0:
        return RGBColor(*default_color)

    roi_rgb = cv2.cvtColor(roi, cv2.COLOR_BGR2RGB)
    pixels = roi_rgb.reshape(-1, 3).astype(np.float32)

    if len(pixels) < 4:
        median = np.median(pixels, axis=0)
        return RGBColor(int(median[0]), int(median[1]), int(median[2]))

    mean_color = np.mean(pixels, axis=0)
    distances = np.linalg.norm(pixels - mean_color, axis=1)
    threshold = max(20.0, float(np.percentile(distances, 70)))
    text_pixels = pixels[distances >= threshold]

    if len(text_pixels) == 0:
        text_pixels = pixels[distances >= np.percentile(distances, 50)]

    median = np.median(text_pixels, axis=0)
    r, g, b = int(median[0]), int(median[1]), int(median[2])

    bg_brightness = float(np.mean(mean_color))
    text_brightness = float(np.mean(median))

    if abs(bg_brightness - text_brightness) < 15:
        if bg_brightness >= 128:
            return RGBColor(0, 0, 0)
        return RGBColor(255, 255, 255)

    return RGBColor(
        max(0, min(255, r)),
        max(0, min(255, g)),
        max(0, min(255, b)),
    )


def _set_shape_no_fill_no_line(shape) -> None:
    """텍스트박스 배경 투명, 테두리 없음."""
    shape.fill.background()
    shape.line.fill.background()


class PptxBuilder:
    """16:9 PPTX 생성기."""

    def __init__(self):
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(PPT_WIDTH_INCH)
        self.presentation.slide_height = Inches(PPT_HEIGHT_INCH)
        self.blank_layout = self.presentation.slide_layouts[6]

    def add_slide(
        self,
        background_image_path: str,
        ocr_boxes: List[OcrBox],
        image_width_px: int,
        image_height_px: int,
        resolved_font: ResolvedFont,
        source_image_for_color: str,
        font_size_scale: float = DEFAULT_FONT_SIZE_SCALE,
        log_callback: Optional[Callable[[str], None]] = None,
    ) -> None:
        """텍스트 제거된 배경 이미지와 OCR 텍스트박스를 슬라이드에 추가합니다."""
        if not os.path.isfile(background_image_path):
            raise FileNotFoundError(
                f"배경 이미지를 찾을 수 없습니다: {background_image_path}"
            )

        slide = self.presentation.slides.add_slide(self.blank_layout)

        slide.shapes.add_picture(
            background_image_path,
            Inches(0),
            Inches(0),
            width=Inches(PPT_WIDTH_INCH),
            height=Inches(PPT_HEIGHT_INCH),
        )

        for box in ocr_boxes:
            ppt_x, ppt_y, ppt_w, ppt_h = px_to_inches(
                box.x,
                box.y,
                box.width,
                box.height,
                image_width_px,
                image_height_px,
            )

            left = Inches(ppt_x)
            top = Inches(ppt_y)
            width = Inches(max(ppt_w, 0.05))
            height = Inches(max(ppt_h, 0.03))

            textbox = slide.shapes.add_textbox(left, top, width, height)
            _set_shape_no_fill_no_line(textbox)

            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.NONE
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0

            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(0)

            run = paragraph.add_run()
            run.text = box.text

            font_size = estimate_font_size_pt(ppt_h, scale=font_size_scale)
            run.font.size = Pt(font_size)
            apply_font_to_run(run, resolved_font)

            try:
                run.font.color.rgb = estimate_text_color(source_image_for_color, box)
            except Exception:
                run.font.color.rgb = RGBColor(0, 0, 0)
                if log_callback:
                    log_callback(
                        f"텍스트 색상 추정 실패, 검정색 적용: '{box.text[:20]}...'"
                    )

    def save(self, output_path: str) -> str:
        """PPTX 파일을 저장합니다."""
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        if not output_path.lower().endswith(".pptx"):
            output_path = output_path + ".pptx"

        try:
            self.presentation.save(output_path)
        except Exception as exc:
            raise RuntimeError(f"PPTX 저장에 실패했습니다: {output_path}") from exc

        return output_path
